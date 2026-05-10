"""Generate the two Capstone PPTX presentations for Hardware Vault.

Run from repo root:  python docs/generate_presentations.py
Outputs:
  docs/presentations/Presentation1_Flujo_Uso.pptx
  docs/presentations/Presentation2_Tecnica.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Output to a separate Downloads folder so deliverables are not committed to the repo.
OUT_DIR = Path.home() / "Downloads" / "Capstone_Entregables"
OUT_DIR.mkdir(parents=True, exist_ok=True)

# ─── Palette (matches app theme) ─────────────────────────────────────────────
BG_DARK = RGBColor(0x0E, 0x1A, 0x12)        # near-black green
BG_CARD = RGBColor(0x16, 0x26, 0x1B)        # surface card
PRIMARY = RGBColor(0x00, 0xC8, 0x53)        # signature green
TEXT_HI = RGBColor(0xF5, 0xF7, 0xF6)        # primary text
TEXT_LO = RGBColor(0xA8, 0xB5, 0xAC)        # muted text
INTEL = RGBColor(0x00, 0x71, 0xC5)
AMD = RGBColor(0xED, 0x1C, 0x24)
NVIDIA = RGBColor(0x76, 0xB9, 0x00)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def new_pres():
    pres = Presentation()
    pres.slide_width = SLIDE_W
    pres.slide_height = SLIDE_H
    return pres


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def set_bg(slide, color=BG_DARK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    fill(bg, color)
    bg.shadow.inherit = False


def textbox(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT_HI,
            font="Calibri", align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        if i > 0:
            p = tf.add_paragraph()
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def header(slide, title, subtitle=None):
    """Standard content-slide header: green accent bar + title (+ optional subtitle)."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.55), Inches(0.12), Inches(0.55))
    fill(bar, PRIMARY)
    textbox(slide, Inches(0.85), Inches(0.45), Inches(11.5), Inches(0.7),
            title, size=28, bold=True, color=TEXT_HI, font="Calibri")
    if subtitle:
        textbox(slide, Inches(0.85), Inches(1.05), Inches(11.5), Inches(0.4),
                subtitle, size=13, color=TEXT_LO, italic=True)


def cover(slide, kicker, title, subtitle, metadata):
    set_bg(slide)
    # diagonal green accent (just a tall rectangle on the right)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(11.8), Inches(0), Inches(1.5), SLIDE_H)
    fill(accent, PRIMARY)

    textbox(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.5),
            kicker, size=14, bold=True, color=PRIMARY)
    textbox(slide, Inches(0.8), Inches(1.9), Inches(10.5), Inches(2.0),
            title, size=54, bold=True, color=TEXT_HI)
    textbox(slide, Inches(0.8), Inches(3.6), Inches(10.5), Inches(0.7),
            subtitle, size=18, color=TEXT_LO, italic=True)
    # metadata stack
    y = Inches(5.0)
    for label, value in metadata:
        textbox(slide, Inches(0.8), y, Inches(2.5), Inches(0.35),
                label, size=12, bold=True, color=TEXT_LO)
        textbox(slide, Inches(3.3), y, Inches(8), Inches(0.35),
                value, size=12, color=TEXT_HI)
        y += Inches(0.4)


def bullets(slide, x, y, w, h, items, *, size=16, gap=0.05):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        # bullet dot
        dot = p.add_run()
        dot.text = "● "
        dot.font.size = Pt(size)
        dot.font.color.rgb = PRIMARY
        dot.font.name = "Calibri"
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.name = "Calibri"
        run.font.color.rgb = TEXT_HI


def card(slide, x, y, w, h, title, body, *, accent=PRIMARY):
    """Dark card with a colored top border, title and body text."""
    base = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    base.adjustments[0] = 0.06
    fill(base, BG_CARD)
    # top accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.07))
    fill(line, accent)
    textbox(slide, x + Inches(0.25), y + Inches(0.25),
            w - Inches(0.5), Inches(0.45),
            title, size=15, bold=True, color=accent)
    textbox(slide, x + Inches(0.25), y + Inches(0.75),
            w - Inches(0.5), h - Inches(1.0),
            body, size=12, color=TEXT_HI)


def stat(slide, x, y, w, value, label, *, color=PRIMARY):
    textbox(slide, x, y, w, Inches(1.6), value,
            size=80, bold=True, color=color, align=PP_ALIGN.CENTER, font="Calibri")
    textbox(slide, x, y + Inches(1.5), w, Inches(0.5), label,
            size=13, color=TEXT_LO, align=PP_ALIGN.CENTER, italic=True)


def table_block(slide, x, y, headers, rows, col_widths_inches,
                row_h=0.4, header_h=0.45):
    """Manual table built from rectangles + text boxes. Avoids native table styling."""
    # header
    cx = x
    for i, htxt in enumerate(headers):
        w = Inches(col_widths_inches[i])
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, y, w, Inches(header_h))
        fill(cell, PRIMARY)
        textbox(slide, cx + Inches(0.1), y + Inches(0.05), w - Inches(0.2),
                Inches(header_h - 0.05), htxt, size=12, bold=True, color=BG_DARK)
        cx += w
    # body
    cy = y + Inches(header_h)
    for r, row in enumerate(rows):
        cx = x
        bg_color = BG_CARD if r % 2 == 0 else BG_DARK
        for i, txt in enumerate(row):
            w = Inches(col_widths_inches[i])
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, w, Inches(row_h))
            fill(cell, bg_color)
            textbox(slide, cx + Inches(0.1), cy + Inches(0.06), w - Inches(0.2),
                    Inches(row_h - 0.05), str(txt), size=11, color=TEXT_HI)
            cx += w
        cy += Inches(row_h)


def screenshot_box(slide, x, y, w, h, label):
    """Placeholder rectangle representing a screenshot."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.04
    fill(box, BG_CARD)
    textbox(slide, x, y + h / 2 - Inches(0.2), w, Inches(0.4),
            f"[Screenshot: {label}]", size=12, color=TEXT_LO,
            italic=True, align=PP_ALIGN.CENTER)


def slide_number(slide, n, total):
    textbox(slide, Inches(12.4), Inches(7.1), Inches(0.8), Inches(0.3),
            f"{n}/{total}", size=10, color=TEXT_LO, align=PP_ALIGN.RIGHT)


# ─── Presentation 1: Flujo / uso de la app ───────────────────────────────────
def build_presentation_1():
    pres = new_pres()
    blank = pres.slide_layouts[6]
    total = 16
    n = 0

    def add():
        nonlocal n
        n += 1
        s = pres.slides.add_slide(blank)
        set_bg(s)
        return s

    # 1. Portada
    s = add()
    cover(s, "PROYECTO CAPSTONE — COMP3402",
          "Hardware Vault",
          "Catálogo móvil de hardware con builds personales y feed de noticias",
          [("Estudiante:", "Alejandro Ruiz"),
           ("Profesor:", "Javier Dastas"),
           ("Curso:", "Ingeniería de Software II — COMP3402"),
           ("Repositorio:", "github.com/AlejandroI28/Final-Project-COMP3402"),
           ("Presentación:", "1 de 2 — Flujo y uso de la aplicación")])
    slide_number(s, n, total)

    # 2. Problema y propósito
    s = add()
    header(s, "Problema y propósito", "El problema fragmentado de armar una PC")
    bullets(s, Inches(0.85), Inches(1.9), Inches(7.5), Inches(4.5), [
        "Las specs de cada componente están dispersas entre páginas de fabricantes.",
        "Las noticias técnicas viven en blogs y canales separados.",
        "Documentar una build personal termina en notas sueltas o spreadsheets.",
        "No hay un solo punto móvil para comparar, leer y registrar.",
    ])
    card(s, Inches(8.7), Inches(1.9), Inches(4.0), Inches(4.5),
         "Propósito",
         "Una app móvil offline-first que unifica:\n\n"
         "• Catálogo de CPUs y GPUs\n"
         "• Feed de noticias del sector\n"
         "• Builds personales con persistencia local en SQLite")
    slide_number(s, n, total)

    # 3. Público objetivo
    s = add()
    header(s, "Público objetivo", "Tres perfiles a los que la app sirve")
    card(s, Inches(0.85), Inches(2.0), Inches(3.9), Inches(4.0),
         "Entusiasta de PC",
         "Conoce hardware, sigue lanzamientos, planea upgrades y disfruta comparar specs.",
         accent=PRIMARY)
    card(s, Inches(4.95), Inches(2.0), Inches(3.9), Inches(4.0),
         "Builder novato",
         "Está armando su primera PC y necesita una vista clara de componentes y rangos de precio.",
         accent=NVIDIA)
    card(s, Inches(9.05), Inches(2.0), Inches(3.6), Inches(4.0),
         "Estudiante",
         "Usa la app como referencia rápida durante cursos de arquitectura de computadoras.",
         accent=INTEL)
    slide_number(s, n, total)

    # 4. Vista general del sistema
    s = add()
    header(s, "Vista general del sistema", "Tres pestañas, un BottomNavigationBar")
    card(s, Inches(0.85), Inches(2.0), Inches(3.9), Inches(4.0),
         "1. Catalog",
         "Búsqueda, filtros (Type/Brand/Series), sort, vista de detalle.\n\n"
         "Empty state inicial — los resultados solo aparecen tras presionar Search.")
    card(s, Inches(4.95), Inches(2.0), Inches(3.9), Inches(4.0),
         "2. News",
         "Feed con card destacado y lista. 8 artículos mock con imágenes Unsplash temáticas y categorías GPU/CPU/Memory/Event.")
    card(s, Inches(9.05), Inches(2.0), Inches(3.6), Inches(4.0),
         "3. My PC",
         "CRUD de builds (CPU, GPU, RAM, Storage, notas).\n\n"
         "Persistencia en SQLite local — sobrevive a kill/reopen.")
    slide_number(s, n, total)

    # 5. Flujo principal de uso
    s = add()
    header(s, "Flujo principal de uso", "Recorrido típico del usuario")
    steps = [
        ("1", "Splash (~1.4 s) → MainShell aterriza en News"),
        ("2", "Usuario navega a Catalog y configura filtros + query"),
        ("3", "Presiona Search → resultados filtrados aparecen"),
        ("4", "Toca un componente → ficha técnica con score bar"),
        ("5", "Va a My PC, presiona +, llena el editor, guarda"),
        ("6", "Reabre la app → su build sigue ahí (SQLite)"),
    ]
    y = Inches(1.9)
    for num, txt in steps:
        circle = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.85), y, Inches(0.7), Inches(0.7))
        fill(circle, PRIMARY)
        textbox(s, Inches(0.85), y + Inches(0.12), Inches(0.7), Inches(0.5),
                num, size=22, bold=True, color=BG_DARK,
                align=PP_ALIGN.CENTER)
        textbox(s, Inches(1.8), y + Inches(0.18), Inches(11), Inches(0.5),
                txt, size=15, color=TEXT_HI)
        y += Inches(0.75)
    slide_number(s, n, total)

    # 6. Pantallas clave
    s = add()
    header(s, "Pantallas clave", "Lo que verá quien use la app")
    titles = ["Splash", "News", "Catalog (empty)",
              "Catalog (results)", "Part Detail", "Build Editor"]
    files = [
        "splash_screen.dart", "news_screen.dart", "catalog_screen.dart",
        "catalog_screen.dart", "part_detail_screen.dart",
        "my_equipment_screen.dart"]
    for i, (t, f) in enumerate(zip(titles, files)):
        col = i % 3
        row = i // 3
        x = Inches(0.85 + col * 4.1)
        y = Inches(1.85 + row * 2.7)
        screenshot_box(s, x, y, Inches(3.9), Inches(2.0), t)
        textbox(s, x, y + Inches(2.05), Inches(3.9), Inches(0.3),
                f"lib/screens/{f}", size=10, color=TEXT_LO, italic=True,
                align=PP_ALIGN.CENTER)
    slide_number(s, n, total)

    # 7. Funcionalidad 1 — Catalog
    s = add()
    header(s, "Funcionalidad 1 — Catálogo con búsqueda intencional",
           "El catálogo no muestra resultados hasta presionar Search")
    bullets(s, Inches(0.85), Inches(1.9), Inches(6.5), Inches(4.5), [
        "Type filter: All / CPU / GPU.",
        "Brand filter dinámico según el Type seleccionado.",
        "Series filter encadenado al Brand.",
        "Sort: Newest / Oldest / Cheapest / Most expensive.",
        "Botón verde Search aplica el query y los filtros activos.",
        "Botón Clear (X) aparece solo cuando hay filtros activos.",
    ])
    screenshot_box(s, Inches(7.7), Inches(1.85), Inches(5.0), Inches(4.7),
                   "Catalog con búsqueda y filtros aplicados")
    slide_number(s, n, total)

    # 8. Funcionalidad 2 — News con search
    s = add()
    header(s, "Funcionalidad 2 — News con búsqueda en vivo",
           "Filtrado por título, resumen, fuente o categoría")
    bullets(s, Inches(0.85), Inches(1.9), Inches(6.5), Inches(4.5), [
        "Search bar superior — filtra en cada keystroke.",
        "Card destacado con la noticia más reciente.",
        "Lista de cards densos con thumbnail Unsplash.",
        "Tag de categoría coloreado (GPU verde, CPU azul, etc.).",
        "Timestamp relativo con paquete timeago en inglés.",
        "Loader (spinner verde) mientras cargan las imágenes.",
    ])
    screenshot_box(s, Inches(7.7), Inches(1.85), Inches(5.0), Inches(4.7),
                   "News con featured + lista")
    slide_number(s, n, total)

    # 9. Funcionalidad 3 — My PC + SQLite
    s = add()
    header(s, "Funcionalidad 3 — My PC con persistencia SQLite",
           "Crear, editar y eliminar builds con DB real")
    bullets(s, Inches(0.85), Inches(1.9), Inches(6.5), Inches(4.5), [
        "BottomSheet editor con todos los campos del build.",
        "Selección de CPU / GPU desde dropdown del catálogo.",
        "Persistencia en tabla pc_builds vía sqflite.",
        "Lista ordenada por updated_at DESC.",
        "Diálogo de confirmación antes de eliminar.",
        "Sobrevive al kill de la app: se restaura desde SQLite.",
    ])
    screenshot_box(s, Inches(7.7), Inches(1.85), Inches(5.0), Inches(4.7),
                   "Build editor BottomSheet")
    slide_number(s, n, total)

    # 10–12. UAT funcionales
    uat_funcional = [
        ("UAT-F-01", "Crear y persistir un nuevo build",
         "Validar que un build creado en My PC se guarda en SQLite y sigue presente al cerrar y reabrir la app.",
         "Crear build → Save → Kill app → Reabrir.",
         "El build sobrevive con todos sus componentes y la nota \"Updated today\"."),
        ("UAT-F-02", "Búsqueda en catálogo respeta el botón Search",
         "Validar que el catálogo permanece vacío hasta que el usuario presiona Search.",
         "Cambiar filtros + escribir query → confirmar empty state → presionar Search.",
         "Solo después de presionar Search aparecen resultados (los Core Ultra 200 en este caso)."),
        ("UAT-F-03", "Eliminar un build con confirmación",
         "Validar el flujo Cancel/Delete del diálogo y la eliminación efectiva en SQLite.",
         "Tocar Delete → Cancel → Tocar Delete → Confirm → Kill app → Reabrir.",
         "Build eliminado de la lista y de la tabla pc_builds; no reaparece tras reabrir."),
    ]
    for uid, title, obj, steps, expected in uat_funcional:
        s = add()
        header(s, f"{uid} — {title}", "UAT Funcional")
        textbox(s, Inches(0.85), Inches(1.9), Inches(2.0), Inches(0.3),
                "Objetivo", size=12, bold=True, color=PRIMARY)
        textbox(s, Inches(0.85), Inches(2.2), Inches(11.5), Inches(0.8),
                obj, size=14, color=TEXT_HI)
        textbox(s, Inches(0.85), Inches(3.2), Inches(2.0), Inches(0.3),
                "Pasos", size=12, bold=True, color=PRIMARY)
        textbox(s, Inches(0.85), Inches(3.5), Inches(11.5), Inches(1.2),
                steps, size=14, color=TEXT_HI)
        textbox(s, Inches(0.85), Inches(4.7), Inches(2.0), Inches(0.3),
                "Resultado", size=12, bold=True, color=PRIMARY)
        textbox(s, Inches(0.85), Inches(5.0), Inches(11.5), Inches(1.2),
                expected, size=14, color=TEXT_HI)
        # estado badge
        badge = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(10.7), Inches(6.4), Inches(2.0), Inches(0.45))
        badge.adjustments[0] = 0.4
        fill(badge, PRIMARY)
        textbox(s, Inches(10.7), Inches(6.45), Inches(2.0), Inches(0.4),
                "Estado: Aprobada", size=12, bold=True, color=BG_DARK,
                align=PP_ALIGN.CENTER)
        slide_number(s, n, total)

    # 13–15. UAT UI
    uat_ui = [
        ("UAT-UI-01", "Empty state del catálogo + Clear condicional",
         "Verificar que el empty state aparece al inicio y el botón Clear (X) solo se renderiza cuando hay filtros activos.",
         "Abrir Catalog → confirmar empty state → escribir → confirmar X aparece → tocar X.",
         "El botón X aparece al iniciar a tipear y desaparece al tocarlo (filtros vuelven a defaults)."),
        ("UAT-UI-02", "Imágenes de News con loader y fallback",
         "Verificar que Image.network usa loadingBuilder (spinner) y errorBuilder (icono) correctamente.",
         "Abrir News con red → confirmar spinner → cargar imagen → cortar red → reabrir.",
         "Con red: spinner → imagen real. Sin red: icono newspaper_rounded en placeholder gris."),
        ("UAT-UI-03", "Splash termina y aterriza en News sin glitch",
         "Validar que el splash simplificado ejecuta su animación, navega correctamente y aterriza en la pestaña News.",
         "Lanzar la app desde frío → cronometrar fade-in → confirmar landing.",
         "Fade-in ~350 ms, total ~1.4 s, sin glitches. News activa por defecto."),
    ]
    for uid, title, obj, steps, expected in uat_ui:
        s = add()
        header(s, f"{uid} — {title}", "UAT UI")
        textbox(s, Inches(0.85), Inches(1.9), Inches(2.0), Inches(0.3),
                "Objetivo", size=12, bold=True, color=PRIMARY)
        textbox(s, Inches(0.85), Inches(2.2), Inches(11.5), Inches(0.8),
                obj, size=14, color=TEXT_HI)
        textbox(s, Inches(0.85), Inches(3.2), Inches(2.0), Inches(0.3),
                "Pasos", size=12, bold=True, color=PRIMARY)
        textbox(s, Inches(0.85), Inches(3.5), Inches(11.5), Inches(1.2),
                steps, size=14, color=TEXT_HI)
        textbox(s, Inches(0.85), Inches(4.7), Inches(2.0), Inches(0.3),
                "Resultado", size=12, bold=True, color=PRIMARY)
        textbox(s, Inches(0.85), Inches(5.0), Inches(11.5), Inches(1.2),
                expected, size=14, color=TEXT_HI)
        badge = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(10.7), Inches(6.4), Inches(2.0), Inches(0.45))
        badge.adjustments[0] = 0.4
        fill(badge, PRIMARY)
        textbox(s, Inches(10.7), Inches(6.45), Inches(2.0), Inches(0.4),
                "Estado: Aprobada", size=12, bold=True, color=BG_DARK,
                align=PP_ALIGN.CENTER)
        slide_number(s, n, total)

    # 16. Conclusión funcional
    s = add()
    header(s, "Conclusión funcional", "Lo que la app entrega y por qué importa")
    stat(s, Inches(0.85), Inches(2.0), Inches(4.0), "3", "secciones funcionales")
    stat(s, Inches(4.85), Inches(2.0), Inches(4.0), "39", "componentes en catálogo")
    stat(s, Inches(8.85), Inches(2.0), Inches(4.0), "6", "UAT aprobadas")
    textbox(s, Inches(0.85), Inches(4.6), Inches(11.5), Inches(2.5),
            "Hardware Vault permite explorar 13 CPUs y 18 GPUs, leer 8 noticias "
            "categorizadas, y mantener un registro persistente de builds personales — "
            "todo offline-first, con UI accesible y consistente en dark mode. La app "
            "respeta el patrón \"búsqueda intencional\" (Search button + Clear "
            "condicional) y demuestra que se puede entregar una experiencia móvil "
            "ligera sin recurrir a servicios en la nube.",
            size=15, color=TEXT_HI)
    slide_number(s, n, total)

    out = OUT_DIR / "Presentation1_Flujo_Uso.pptx"
    pres.save(out)
    print(f"OK - {out}")


# ─── Presentation 2: Aspectos técnicos ───────────────────────────────────────
def build_presentation_2():
    pres = new_pres()
    blank = pres.slide_layouts[6]
    total = 16
    n = 0

    def add():
        nonlocal n
        n += 1
        s = pres.slides.add_slide(blank)
        set_bg(s)
        return s

    # 1. Portada
    s = add()
    cover(s, "PROYECTO CAPSTONE — COMP3402",
          "Hardware Vault",
          "Aspectos técnicos: arquitectura, base de datos, pruebas y DevOps",
          [("Estudiante:", "Alejandro Ruiz"),
           ("Profesor:", "Javier Dastas"),
           ("Curso:", "Ingeniería de Software II — COMP3402"),
           ("Repositorio:", "github.com/AlejandroI28/Final-Project-COMP3402"),
           ("Presentación:", "2 de 2 — Aspectos técnicos")])
    slide_number(s, n, total)

    # 2. Tecnologías utilizadas
    s = add()
    header(s, "Tecnologías utilizadas", "Stack móvil sin backend")
    table_block(s, Inches(0.85), Inches(1.9),
                ["Categoría", "Tecnología", "Versión / nota"],
                [
                    ["Lenguaje", "Dart", "3.0+"],
                    ["Framework", "Flutter", "3.x — Material 3, dark mode"],
                    ["Estado", "provider", "6.1.1 — ChangeNotifier"],
                    ["Persistencia", "sqflite + path", "2.3.0 — SQLite local"],
                    ["Tipografía", "google_fonts", "Space Grotesk"],
                    ["Fechas", "timeago", "3.6.0"],
                    ["Imágenes red", "Image.network nativo", "Unsplash CDN"],
                    ["Imágenes locales", "Image.asset", "assets/images/*"],
                    ["Testing", "flutter_test", "14 unit tests pasando"],
                    ["Plataforma", "Android", "emulador Pixel 3a API 34"],
                ],
                col_widths_inches=[2.4, 3.6, 5.5], row_h=0.42)
    slide_number(s, n, total)

    # 3. Arquitectura del sistema
    s = add()
    header(s, "Arquitectura del sistema", "Capas — Presentation / State / Data / Persistence")
    layers = [
        ("Presentation", "Splash, MainShell, Catalog, News, MyEquipment, PartDetail, _BuildEditor + widgets compartidos (BrandBadge, ScoreBar, EmptyState).", PRIMARY),
        ("State", "AppState (extends ChangeNotifier) expuesto vía Provider. Maneja filtros del catálogo, máquina de estados de búsqueda y lista de builds.", NVIDIA),
        ("Data", "Models (CPU, GPU, NewsArticle, PCBuild), mock_data.dart (catálogo read-only) y DatabaseHelper singleton para SQLite.", INTEL),
        ("Persistence", "SQLite local — archivo hardware_vault.db con tabla pc_builds.", AMD),
    ]
    y = Inches(1.95)
    for name, desc, color in layers:
        bar = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.85), y, Inches(0.15), Inches(1.05))
        fill(bar, color)
        textbox(s, Inches(1.1), y, Inches(2.5), Inches(0.5),
                name, size=18, bold=True, color=color)
        textbox(s, Inches(1.1), y + Inches(0.5), Inches(11.5), Inches(0.6),
                desc, size=12, color=TEXT_HI)
        y += Inches(1.2)
    slide_number(s, n, total)

    # 4. Componentes del sistema
    s = add()
    header(s, "Componentes del sistema",
           "Frontend, estado, datos — sin backend ni servicios externos")
    card(s, Inches(0.85), Inches(2.0), Inches(3.9), Inches(4.5),
         "Frontend",
         "Widgets Flutter Material 3 con dark mode exclusivo.\n\n"
         "Pantallas: Splash, MainShell, Catalog, News, MyEquipment, PartDetail.\n\n"
         "Widgets reutilizables en shared_widgets.dart.")
    card(s, Inches(4.95), Inches(2.0), Inches(3.9), Inches(4.5),
         "Estado",
         "Provider + ChangeNotifier.\n\n"
         "AppState centraliza:\n"
         "• Filtros del catálogo (search, brand, series, sort)\n"
         "• Máquina de hasSearched / clearCatalogFilters\n"
         "• Lista in-memory de builds")
    card(s, Inches(9.05), Inches(2.0), Inches(3.6), Inches(4.5),
         "Datos",
         "Catálogo read-only en mock_data.dart.\n\n"
         "DatabaseHelper singleton expone CRUD sobre SQLite.\n\n"
         "Modelos inmutables (CPU, GPU, NewsArticle) y mutable (PCBuild).")
    slide_number(s, n, total)

    # 5. Base de datos
    s = add()
    header(s, "Base de datos", "Tabla pc_builds — DDL aplicada por DatabaseHelper")
    table_block(s, Inches(0.85), Inches(1.85),
                ["Columna", "Tipo", "Default", "Descripción"],
                [
                    ["id", "TEXT PK", "—", "build_<millisSinceEpoch>"],
                    ["name", "TEXT NOT NULL", "—", "Nombre del build"],
                    ["cpu_id", "TEXT", "—", "FK lógica a CPU.id"],
                    ["gpu_id", "TEXT", "—", "FK lógica a GPU.id"],
                    ["ram_gb", "INTEGER", "16", "8/16/32/64/128"],
                    ["ram_type", "TEXT", "DDR5", "DDR4 / DDR5"],
                    ["storage_gb", "INTEGER", "1000", "256–4000"],
                    ["storage_type", "TEXT", "NVMe SSD", "NVMe / SATA / HDD"],
                    ["psu_watts", "TEXT", "—", "Opcional"],
                    ["case_model", "TEXT", "—", "Opcional"],
                    ["notes", "TEXT", "''", "Notas libres"],
                    ["created_at", "TEXT", "—", "ISO-8601"],
                    ["updated_at", "TEXT", "—", "ISO-8601 (sort key)"],
                ],
                col_widths_inches=[2.0, 2.0, 1.5, 6.0], row_h=0.34)
    slide_number(s, n, total)

    # 6. Flujo técnico de datos
    s = add()
    header(s, "Flujo técnico de datos", "Cómo entra, se procesa y se guarda")
    steps = [
        "Usuario abre _BuildEditor → llena formulario.",
        "Tap Save → AppState.saveBuild(build).",
        "AppState llama DatabaseHelper.instance.upsertBuild(build).",
        "DatabaseHelper convierte PCBuild → Map<String, Object?> con _toRow().",
        "INSERT OR REPLACE INTO pc_builds VALUES (...).",
        "AppState actualiza la lista in-memory + notifyListeners().",
        "MyEquipmentScreen rebuild — el card aparece con updated_at = now.",
    ]
    y = Inches(1.95)
    for i, step in enumerate(steps, 1):
        circle = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.85), y, Inches(0.55), Inches(0.55))
        fill(circle, PRIMARY)
        textbox(s, Inches(0.85), y + Inches(0.07), Inches(0.55), Inches(0.4),
                str(i), size=18, bold=True, color=BG_DARK,
                align=PP_ALIGN.CENTER)
        textbox(s, Inches(1.6), y + Inches(0.13), Inches(11.2), Inches(0.5),
                step, size=14, color=TEXT_HI)
        y += Inches(0.65)
    slide_number(s, n, total)

    # 7. Organización del repositorio
    s = add()
    header(s, "Organización del repositorio", "Estructura por capas")
    code = (
        "hardware_vault/\n"
        "├── lib/\n"
        "│   ├── main.dart                # Bootstrap + MainShell\n"
        "│   ├── theme/app_theme.dart\n"
        "│   ├── models/models.dart       # CPU, GPU, NewsArticle, PCBuild\n"
        "│   ├── data/\n"
        "│   │   ├── mock_data.dart       # Catálogo read-only\n"
        "│   │   └── database.dart        # SQLite helper (sqflite)\n"
        "│   ├── providers/app_state.dart # ChangeNotifier global\n"
        "│   ├── screens/                 # 5 pantallas\n"
        "│   └── widgets/shared_widgets.dart\n"
        "├── test/                        # 14 unit tests\n"
        "├── assets/images/               # Logos + product photos\n"
        "├── docs/                        # Diagramas, UAT, SCRUM, PDF, PPTX\n"
        "├── pubspec.yaml\n"
        "└── README.md")
    bg = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.85), Inches(1.9), Inches(11.5), Inches(5.0))
    bg.adjustments[0] = 0.03
    fill(bg, BG_CARD)
    textbox(s, Inches(1.1), Inches(2.1), Inches(11.0), Inches(4.6),
            code, size=13, color=TEXT_HI, font="Consolas")
    slide_number(s, n, total)

    # 8. README.md
    s = add()
    header(s, "README.md", "Pieza central de documentación")
    bullets(s, Inches(0.85), Inches(1.95), Inches(11.5), Inches(5.0), [
        "Título, integrante, profesor, curso y URL del repo.",
        "Descripción del problema y objetivos del sistema.",
        "Tabla de tecnologías + diagrama de arquitectura por capas.",
        "Módulos / componentes con enlaces clicables a archivos lib/.",
        "Instrucciones para ejecutar (clone, pub get, run, test).",
        "Estructura del repositorio y descripción de la base de datos.",
        "Evidencia de SCRUM (6 sprints) y de testing (14 unit + 6 UAT).",
        "Estado actual, mejoras futuras, enlaces a videos y PPTX.",
        "Capturas clave referenciando archivos fuente.",
        "Reflexión individual + reflexión grupal.",
    ], size=14)
    slide_number(s, n, total)

    # 9. Evidencia de SCRUM
    s = add()
    header(s, "Evidencia de SCRUM", "6 sprints — backlog en docs/scrum_evidence.md")
    table_block(s, Inches(0.85), Inches(1.95),
                ["Sprint", "Foco", "Entregables clave"],
                [
                    ["1", "Scaffolding", "Bootstrap, theme, modelos, mock data, navegación"],
                    ["2", "Catálogo + News base", "CatalogScreen + filtros, NewsScreen + feed"],
                    ["3", "My PC", "CRUD de builds, BottomSheet editor"],
                    ["4", "i18n + UX polish", "Traducción a inglés, splash simple, imágenes News"],
                    ["5", "UX catálogo avanzado", "Search button, Clear condicional, modo All"],
                    ["6", "Hardening + entregables", "SQLite, 14 tests, diagramas, UAT, PDF, PPTX"],
                ],
                col_widths_inches=[1.0, 2.5, 8.0], row_h=0.5)
    slide_number(s, n, total)

    # 10. Evidencia de GitHub
    s = add()
    header(s, "Evidencia de GitHub", "Commits, ramas y colaboración")
    bullets(s, Inches(0.85), Inches(1.95), Inches(11.5), Inches(5.0), [
        "Repositorio público en github.com/AlejandroI28/Final-Project-COMP3402",
        "Commits frecuentes con mensajes descriptivos (un commit por feature/fix).",
        "Rama main siempre estable y compilando.",
        "Carpetas organizadas por capa (lib/, test/, docs/, assets/).",
        "README.md como landing page con todos los enlaces internos.",
        "docs/ contiene diagramas, UAT, SCRUM, PDF formal y PPTX.",
    ], size=14)
    slide_number(s, n, total)

    # 11. Pruebas unitarias
    s = add()
    header(s, "Pruebas unitarias", "flutter test — 14 tests pasando")
    table_block(s, Inches(0.85), Inches(1.95),
                ["Grupo", "# tests", "Cubre"],
                [
                    ["Catalog filtering — CPUs", "3", "All sin filtro / brand / search case-insensitive"],
                    ["Catalog filtering — GPUs", "2", "Brand Nvidia / sort price_asc"],
                    ["Catalog search state", "5", "hasSearched / runSearch / clear / hasActive / reset tab"],
                    ["Catalog \"All\" tab", "2", "availableBrands / setBrandFilter sincroniza"],
                    ["AppState notifications", "2", "notifyListeners en setSearch / setSortBy"],
                ],
                col_widths_inches=[3.5, 1.2, 6.8], row_h=0.5)
    textbox(s, Inches(0.85), Inches(6.5), Inches(11.5), Inches(0.4),
            "Output: 00:00 +14: All tests passed!",
            size=13, color=PRIMARY, italic=True, font="Consolas")
    slide_number(s, n, total)

    # 12. UAT funcionales
    s = add()
    header(s, "UAT funcionales", "3 casos documentados con la estructura del rubric")
    table_block(s, Inches(0.85), Inches(1.95),
                ["ID", "Caso", "Estado"],
                [
                    ["UAT-F-01", "Crear y persistir un nuevo build (verifica SQLite tras kill/reopen)", "Aprobada"],
                    ["UAT-F-02", "Búsqueda en catálogo respeta el botón Search", "Aprobada"],
                    ["UAT-F-03", "Eliminar un build con confirmación (diálogo + DELETE en SQLite)", "Aprobada"],
                ],
                col_widths_inches=[1.5, 8.5, 1.5], row_h=0.65)
    textbox(s, Inches(0.85), Inches(5.5), Inches(11.5), Inches(0.4),
            "Detalle completo en docs/uat_tests.md",
            size=13, color=TEXT_LO, italic=True)
    slide_number(s, n, total)

    # 13. UAT UI
    s = add()
    header(s, "UAT de UI", "3 casos visuales / interacción")
    table_block(s, Inches(0.85), Inches(1.95),
                ["ID", "Caso", "Estado"],
                [
                    ["UAT-UI-01", "Empty state del catálogo + aparición condicional del botón Clear", "Aprobada"],
                    ["UAT-UI-02", "Imágenes de News con loader (spinner) y fallback (icono)", "Aprobada"],
                    ["UAT-UI-03", "Splash termina y aterriza en News sin glitch", "Aprobada"],
                ],
                col_widths_inches=[1.5, 8.5, 1.5], row_h=0.65)
    textbox(s, Inches(0.85), Inches(5.5), Inches(11.5), Inches(0.4),
            "Cada UAT incluye objetivo, precondiciones, pasos, resultado esperado/obtenido, estado y evidencia.",
            size=13, color=TEXT_LO, italic=True)
    slide_number(s, n, total)

    # 14. Retos técnicos
    s = add()
    header(s, "Retos técnicos", "Problemas encontrados y cómo se resolvieron")
    card(s, Inches(0.85), Inches(2.0), Inches(3.9), Inches(4.5),
         "Splash pesado",
         "El splash original usaba CustomPainter + glow loop + 3.5 s.\n\n"
         "Causaba lag en emuladores low-end.\n\n"
         "Solución: simplificación a fade-in con AnimatedOpacity, ~1.4 s total.")
    card(s, Inches(4.95), Inches(2.0), Inches(3.9), Inches(4.5),
         "Persistencia débil",
         "SharedPreferences técnicamente persiste, pero no satisface el requisito de \"base de datos\" del rubric.\n\n"
         "Solución: migración a SQLite (sqflite). API de AppState intacta; UI no cambió.",
         accent=NVIDIA)
    card(s, Inches(9.05), Inches(2.0), Inches(3.6), Inches(4.5),
         "Sync TextField↔State",
         "Al limpiar filtros desde el botón Clear, el TextField mantenía el texto.\n\n"
         "Solución: convertir _SearchBar a Stateful con TextEditingController + listener al AppState.",
         accent=AMD)
    slide_number(s, n, total)

    # 15. DevOps / CI-CD
    s = add()
    header(s, "Aspectos DevOps / CI-CD", "Lo implementado y los próximos pasos")
    card(s, Inches(0.85), Inches(2.0), Inches(5.9), Inches(4.5),
         "Implementado",
         "• Repositorio Git en GitHub con commits frecuentes.\n"
         "• Carpetas organizadas por capa.\n"
         "• Test suite local: flutter test (14 tests).\n"
         "• Documentación versionada en docs/ junto al código.\n"
         "• .gitignore configurado para no commitear build/, .dart_tool/, etc.")
    card(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(4.5),
         "Próximos pasos (CI/CD)",
         "• GitHub Actions: flutter analyze + flutter test en cada push.\n"
         "• Job adicional: flutter build apk --release como artifact.\n"
         "• Branch protection en main exigiendo CI verde.\n"
         "• Codecov o similar para tracking de coverage.\n"
         "• Distribución de APK vía GitHub Releases.",
         accent=NVIDIA)
    slide_number(s, n, total)

    # 16. Conclusión técnica
    s = add()
    header(s, "Conclusión técnica", "Lo que el proyecto demuestra")
    stat(s, Inches(0.85), Inches(1.9), Inches(4.0), "14", "unit tests pasando")
    stat(s, Inches(4.85), Inches(1.9), Inches(4.0), "6", "sprints SCRUM")
    stat(s, Inches(8.85), Inches(1.9), Inches(4.0), "1", "tabla SQLite")
    textbox(s, Inches(0.85), Inches(4.5), Inches(11.5), Inches(2.5),
            "El proyecto demuestra que se puede entregar una app móvil funcional con "
            "persistencia real en base de datos, pruebas automáticas, y documentación "
            "rigurosa, dentro del alcance de un curso. La arquitectura por capas y el "
            "uso de Provider + ChangeNotifier hicieron posible un set de unit tests "
            "deterministas sobre la lógica de negocio sin tocar la UI ni la DB. La "
            "decisión de migrar a SQLite habilitó un diagrama ER real y un diccionario "
            "de datos defendible frente al rubric del Capstone.",
            size=14, color=TEXT_HI)
    slide_number(s, n, total)

    out = OUT_DIR / "Presentation2_Tecnica.pptx"
    pres.save(out)
    print(f"OK - {out}")


if __name__ == "__main__":
    build_presentation_1()
    build_presentation_2()
